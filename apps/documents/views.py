import json
import logging
import os

from django.conf import settings
from django.core.files.base import ContentFile
from django.http import FileResponse, Http404, HttpResponse
from django.utils.html import escape as html_escape
from rest_framework import viewsets, permissions, status
from rest_framework.decorators import action
from rest_framework.response import Response
from django_filters.rest_framework import DjangoFilterBackend
from rest_framework.filters import SearchFilter

from apps.accounts.permissions import CanEditProject, org_queryset_filter
from .models import Document
from .serializers import DocumentSerializer

logger = logging.getLogger(__name__)

# MIME type → OnlyOffice fileType mapping (content-based, not filename-based)
_MIME_TO_EXT: dict[str, str] = {
    'application/msword':                                                          'doc',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document':    'docx',
    'application/vnd.oasis.opendocument.text':                                    'odt',
    'application/rtf':                                                             'rtf',
    'text/rtf':                                                                    'rtf',
    'text/plain':                                                                  'txt',
    'application/pdf':                                                             'pdf',
    'application/vnd.ms-excel':                                                    'xls',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':           'xlsx',
    'application/vnd.oasis.opendocument.spreadsheet':                              'ods',
    'text/csv':                                                                    'csv',
    'application/csv':                                                             'csv',
    'application/vnd.ms-powerpoint':                                               'ppt',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation':   'pptx',
    'application/vnd.oasis.opendocument.presentation':                             'odp',
}


def _doc_ext(doc: 'Document') -> str:
    """
    Return the correct OnlyOffice fileType for a Document.
    Prefers the stored MIME type (content-based) over the filename extension so
    that mis-named files (e.g. a PDF saved as .docx) still open correctly.
    """
    if doc.mime_type:
        from_mime = _MIME_TO_EXT.get(doc.mime_type.split(';')[0].strip().lower())
        if from_mime:
            return from_mime
    raw_name = doc.file.name.rsplit('/', 1)[-1]
    return raw_name.rsplit('.', 1)[-1].lower() if '.' in raw_name else 'docx'


def _doc_type(ext: str) -> str:
    if ext in ('doc', 'docx', 'docm', 'dot', 'dotx', 'odt', 'fodt', 'rtf', 'txt', 'pdf', 'epub', 'fb2'):
        return 'word'
    if ext in ('xls', 'xlsx', 'xlsm', 'xlt', 'xltx', 'ods', 'fods', 'csv'):
        return 'cell'
    return 'slide'


class DocumentViewSet(viewsets.ModelViewSet):
    serializer_class = DocumentSerializer
    filter_backends = [DjangoFilterBackend, SearchFilter]
    filterset_fields = ['project', 'folder', 'category', 'ai_processed']
    search_fields = ['title']

    def get_queryset(self):
        from django.db.models import Q
        from apps.accounts.permissions import get_shared_project_ids, get_approved_area_ids
        user = self.request.user
        base_qs = Document.objects.select_related('project__organisation', 'uploaded_by', 'folder')
        own_qs = org_queryset_filter(user, base_qs, org_field='project__organisation')
        if user.is_superadmin or user.role == 'PDDE_VIEWER':
            return own_qs
        shared_project_ids = get_shared_project_ids(user)
        approved_area_ids  = get_approved_area_ids(user)
        if not shared_project_ids and not approved_area_ids:
            return own_qs
        from apps.survey_projects.models import SurveyArea, ProjectLayerFolder
        from collections import deque
        approved_folder_ids: list[int] = []
        for area in SurveyArea.objects.filter(id__in=approved_area_ids).select_related('folder'):
            if area.folder_id:
                queue: deque = deque([area.folder_id])
                while queue:
                    cur = queue.popleft()
                    approved_folder_ids.append(cur)
                    for cid in ProjectLayerFolder.objects.filter(parent_id=cur).values_list('id', flat=True):
                        queue.append(cid)
        extra_q = Q(project_id__in=shared_project_ids)
        if approved_folder_ids:
            extra_q |= Q(folder_id__in=approved_folder_ids)
        return (own_qs | base_qs.filter(extra_q)).distinct()

    def get_permissions(self):
        if self.action == 'embed':
            return [permissions.AllowAny()]
        if self.action == 'onlyoffice_callback':
            return [permissions.AllowAny()]
        if self.action in ['create', 'update', 'partial_update']:
            return [CanEditProject()]
        # destroy: org-level check is enforced in perform_destroy, not via
        # CanEditProject, because CanEditProject.has_permission requires
        # can_forward (SDO/SURVEYOR only) and has_object_permission falls
        # through to a status/organisation_id check that Document doesn't expose.
        return [permissions.IsAuthenticated()]

    def perform_destroy(self, instance):
        from rest_framework.exceptions import PermissionDenied
        from apps.accounts.models import User as _User
        from apps.survey_projects.models import SurveyProject

        user = self.request.user

        if user.is_superadmin:
            instance.delete()
            return

        # Must be same organisation
        if instance.project_id and instance.project.organisation_id != user.organisation_id:
            raise PermissionDenied('You can only delete documents belonging to your organisation.')

        # Admin roles (DEO_ADMIN, CEO_ADMIN, ADEO_ADMIN) can delete regardless of project status
        if user.role in _User.ADMIN_ROLES:
            instance.delete()
            return

        # SDO / SURVEYOR — only allowed when project is DRAFT or RETURNED
        if instance.project_id:
            proj_status = instance.project.status
            if proj_status not in (SurveyProject.DRAFT, SurveyProject.RETURNED):
                raise PermissionDenied(
                    f'Documents can only be deleted while the project is in DRAFT or RETURNED '
                    f'status (current: {proj_status}).'
                )

        instance.delete()

    def perform_create(self, serializer):
        file = self.request.FILES.get('file')
        mime_type = ''
        file_size = 0
        if file:
            content_type = getattr(file, 'content_type', '')
            if content_type and content_type != 'application/octet-stream':
                mime_type = content_type
            try:
                import magic
                detected = magic.from_buffer(file.read(2048), mime=True)
                if detected:
                    mime_type = detected
                file.seek(0)
            except Exception:
                pass
            if not mime_type or mime_type == 'application/octet-stream':
                ext = file.name.split('.')[-1].lower() if '.' in file.name else ''
                mime_map = {
                    'doc': 'application/msword',
                    'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    'odt': 'application/vnd.oasis.opendocument.text',
                    'rtf': 'application/rtf',
                    'txt': 'text/plain',
                    'pdf': 'application/pdf',
                    'xls': 'application/vnd.ms-excel',
                    'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                    'ods': 'application/vnd.oasis.opendocument.spreadsheet',
                    'csv': 'text/csv',
                    'ppt': 'application/vnd.ms-powerpoint',
                    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    'odp': 'application/vnd.oasis.opendocument.presentation',
                }
                mime_type = mime_map.get(ext, 'application/octet-stream')

            # Embed Living Provenance DNA watermark
            try:
                from apps.core.watermark import embed_watermark
                from django.core.files.base import ContentFile
                file.seek(0)
                file_bytes = file.read()
                project = serializer.validated_data.get('project')
                metadata = {
                    "project_id": project.id if project else None,
                    "project_number": project.project_number if project else None,
                    "uploaded_by": self.request.user.username,
                }
                watermarked_bytes = embed_watermark(file_bytes, file.name, mime_type, metadata)
                file_size = len(watermarked_bytes)
                # Overwrite standard file field in validated data
                serializer.validated_data['file'] = ContentFile(watermarked_bytes, name=file.name)
            except Exception as wexc:
                import logging
                logger = logging.getLogger(__name__)
                logger.error(f"Failed to watermark uploaded document: {wexc}")
                file_size = file.size
        serializer.save(
            uploaded_by=self.request.user,
            mime_type=mime_type,
            file_size=file_size,
        )

    @action(detail=True, methods=['get'], url_path='editor-config')
    def editor_config(self, request, pk=None):
        """Return JWT-signed OnlyOffice DocsAPI config for this document."""
        import jwt

        doc = self.get_object()
        secret = getattr(settings, 'ONLYOFFICE_JWT_SECRET', '')
        internal_base = getattr(settings, 'ONLYOFFICE_INTERNAL_BASE_URL', '').rstrip('/')

        if internal_base:
            doc_url = f"{internal_base}{doc.file.url}"
            callback_url = f"{internal_base}/api/documents/{doc.id}/onlyoffice-callback/"
        else:
            doc_url = request.build_absolute_uri(doc.file.url)
            callback_url = request.build_absolute_uri(f'/api/documents/{doc.id}/onlyoffice-callback/')

        ext      = _doc_ext(doc)
        doc_type = _doc_type(ext)

        can_edit = (
            request.user.is_superadmin
            or request.user.organisation_id == doc.project.organisation_id
        )

        title = doc.title if '.' in doc.title else f"{doc.title}.{ext}"

        config = {
            "document": {
                "fileType": ext,
                "key": f"doc-{doc.id}-v{doc.version}",
                "title": title,
                "url": doc_url,
                "permissions": {
                    "edit": can_edit,
                    "download": True,
                    "print": True,
                    "review": can_edit,
                },
            },
            "documentType": doc_type,
            "editorConfig": {
                "callbackUrl": callback_url,
                "lang": "en",
                "mode": "edit" if can_edit else "view",
                "user": {
                    "id": str(request.user.id),
                    "name": request.user.get_full_name() or request.user.username,
                },
                "customization": {
                    "autosave": True,
                    # forcesave=True so Ctrl+S / programmatic force-save persists to storage
                    # immediately (the auto-save-on-close callback is a known weak point).
                    "forcesave": True,
                    "chat": False,
                    "compactHeader": True,
                },
            },
        }

        if secret:
            token = jwt.encode(config, secret, algorithm='HS256')
            if isinstance(token, bytes):
                token = token.decode('utf-8')
            config['token'] = token

        return Response(config)

    @action(detail=True, methods=['get'], url_path='embed',
            permission_classes=[permissions.AllowAny])
    def embed(self, request, pk=None):
        """
        GET /api/documents/{id}/embed/?token=<jwt>

        Standalone HTML page that opens a document in the OnlyOffice editor.
        Authentication is via JWT passed in the `token` query parameter because
        browser tabs opened with window.open() cannot send a Bearer header.

        The OnlyOffice server fetches the document from the internal Docker URL;
        the browser never navigates to that URL directly.
        """
        import jwt as _jwt
        from rest_framework_simplejwt.tokens import UntypedToken
        from rest_framework_simplejwt.exceptions import InvalidToken, TokenError
        from django.contrib.auth import get_user_model

        # ── Authenticate via ?token= query parameter ──────────────────────
        token_param = request.query_params.get('token', '')
        if not token_param:
            return HttpResponse(
                '<h2 style="font-family:sans-serif;color:#c00;padding:24px">'
                'Missing ?token= parameter. '
                'Open documents via the RakshaGIS interface.</h2>',
                content_type='text/html', status=401,
            )
        try:
            validated = UntypedToken(token_param)
            User = get_user_model()
            user = User.objects.get(id=validated['user_id'])
        except (InvalidToken, TokenError, Exception):
            return HttpResponse(
                '<h2 style="font-family:sans-serif;color:#c00;padding:24px">'
                'Session expired or invalid token. '
                'Please return to RakshaGIS and try again.</h2>',
                content_type='text/html', status=401,
            )
        request.user = user   # attach for permission checks below

        doc = self.get_object()
        secret = getattr(settings, 'ONLYOFFICE_JWT_SECRET', '')
        internal_base = getattr(settings, 'ONLYOFFICE_INTERNAL_BASE_URL', '').rstrip('/')

        # document.url  → OnlyOffice server fetches this (internal Docker URL)
        # callbackUrl   → OnlyOffice server POSTs save events here (internal Docker URL)
        if internal_base:
            doc_url = f"{internal_base}{doc.file.url}"
            callback_url = f"{internal_base}/api/documents/{doc.id}/onlyoffice-callback/"
        else:
            doc_url = request.build_absolute_uri(doc.file.url)
            callback_url = request.build_absolute_uri(f'/api/documents/{doc.id}/onlyoffice-callback/')

        ext      = _doc_ext(doc)
        doc_type = _doc_type(ext)

        can_edit = (
            request.user.is_superadmin
            or request.user.organisation_id == doc.project.organisation_id
        )
        title = doc.title if '.' in doc.title else f"{doc.title}.{ext}"

        config = {
            "document": {
                "fileType": ext,
                "key": f"doc-{doc.id}-v{doc.version}",
                "title": title,
                "url": doc_url,
                "permissions": {
                    "edit": can_edit, "download": True,
                    "print": True, "review": can_edit,
                },
            },
            "documentType": doc_type,
            "editorConfig": {
                "callbackUrl": callback_url,
                "lang": "en",
                "mode": "edit" if can_edit else "view",
                "user": {
                    "id": str(request.user.id),
                    "name": request.user.get_full_name() or request.user.username,
                },
                "customization": {
                    "autosave": True, "forcesave": True,
                    "chat": False, "compactHeader": True,
                    "goback": {"url": request.META.get('HTTP_REFERER', '/documents')},
                },
            },
        }

        if secret:
            token = _jwt.encode(config, secret, algorithm='HS256')
            if isinstance(token, bytes):
                token = token.decode('utf-8')
            config['token'] = token

        # Embed JSON safely inside a <script> tag: escape only the characters
        # that could break out of the script context. Do NOT html-escape the
        # whole blob — that would corrupt the JSON (quotes → &quot; etc).
        config_json = (
            json.dumps(config)
            .replace('<', '\\u003c')
            .replace('>', '\\u003e')
            .replace('&', '\\u0026')
        )
        # title_safe goes into HTML text content → Django's escape (single arg)
        title_safe = html_escape(title)

        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8"/>
  <meta name="viewport" content="width=device-width,initial-scale=1"/>
  <title>{title_safe} — RakshaGIS</title>
  <style>
    * {{ margin:0; padding:0; box-sizing:border-box; }}
    body {{ background:#1a1a2e; font-family:Arial,sans-serif; height:100vh; display:flex; flex-direction:column; }}
    #header {{
      background:#0d2b5e; color:#fff; padding:8px 16px;
      display:flex; align-items:center; gap:12px; flex-shrink:0;
      font-size:14px; border-bottom:2px solid #1e5091;
    }}
    #header strong {{ font-size:15px; }}
    #editor {{ flex:1; }}
    .loading {{
      color:#aaa; text-align:center; padding:60px;
      font-size:16px;
    }}
    .error {{
      color:#ff6b6b; background:#1e1e1e; padding:24px;
      border:1px solid #ff4d4f; border-radius:4px; margin:20px;
    }}
  </style>
</head>
<body>
  <div id="header">
    <span>🗂</span>
    <strong>{title_safe}</strong>
    <span style="color:#90b8d8;font-size:12px">RakshaGIS — DGDE Survey Platform</span>
  </div>
  <div id="editor"><div class="loading">Loading document editor…</div></div>
  <script>
    (function() {{
      var config = {config_json};
      config.width  = '100%';
      config.height = '100%';
      config.events = {{
        onError: function(e) {{
          document.getElementById('editor').innerHTML =
            '<div class="error"><b>Editor error:</b> ' +
            (e && e.data ? e.data : 'Unknown error') + '</div>';
        }}
      }};

      function initEditor() {{
        try {{
          new DocsAPI.DocEditor('editor', config);
        }} catch(err) {{
          showError("Failed to initialize editor: " + err.message);
        }}
      }}

      function showError(msg) {{
        document.getElementById('editor').innerHTML =
          '<div class="error"><b>OnlyOffice Unavailable:</b> ' + msg +
          '<br/><br/><button onclick="window.location.reload()" style="padding:6px 12px;background:#ff4d4f;color:#fff;border:none;border-radius:4px;cursor:pointer">Retry Loading Page</button></div>';
      }}

      // Load OnlyOffice API dynamically to handle startup delay/errors gracefully
      var script = document.createElement('script');
      script.src = "/onlyoffice/web-apps/apps/api/documents/api.js";
      script.onload = function() {{
        if (typeof DocsAPI !== 'undefined') {{
          initEditor();
        }} else {{
          showError("OnlyOffice API script loaded, but DocsAPI is undefined. The service may still be starting up.");
        }}
      }};
      script.onerror = function() {{
        showError("Could not load OnlyOffice API script. Make sure the OnlyOffice container is running and healthy, or reload the page in a few seconds.");
      }};
      document.body.appendChild(script);
    }})();
  </script>
</body>
</html>"""
        return HttpResponse(html, content_type='text/html')

    @action(detail=True, methods=['post'], url_path='onlyoffice-callback',
            permission_classes=[permissions.AllowAny], authentication_classes=[])
    def onlyoffice_callback(self, request, pk=None):
        """OnlyOffice calls this endpoint when the document is saved."""
        import jwt
        import urllib.request

        secret = getattr(settings, 'ONLYOFFICE_JWT_SECRET', '')

        if secret:
            # OnlyOffice normally includes the callback JWT in the POST body
            # as `token`. Prefer that value to avoid DRF/SimpleJWT attempting
            # to authenticate using other Bearer tokens in the Authorization
            # header (which are unrelated and may cause 401s).
            token = request.data.get('token', '') or ''
            if not token:
                auth_header = request.headers.get('Authorization', '')
                if auth_header.startswith('Bearer '):
                    token = auth_header[7:]
            if token:
                try:
                    jwt.decode(token, secret, algorithms=['HS256'])
                except Exception:
                    return Response({'error': 1}, status=403)
            else:
                return Response({'error': 1}, status=403)

        body = request.data
        oo_status = body.get('status')
        # OnlyOffice status codes: 1=editing, 2=ready-to-save (all editors closed),
        # 3=save error, 4=closed-no-changes, 6=force-save while editing, 7=force-save error.
        if oo_status in (3, 7):
            logger.error("OnlyOffice reported a save error (status %s) for doc %s: %s",
                         oo_status, pk, body.get('url') or body)

        if oo_status in (2, 6):
            download_url = body.get('url', '')
            doc = Document.objects.filter(pk=pk).first()
            if not doc:
                logger.error("OnlyOffice callback: document %s not found", pk)
                return Response({'error': 1})
            if not download_url:
                logger.error("OnlyOffice callback for doc %s: status %s but no download URL "
                             "in payload — edits cannot be retrieved", pk, oo_status)
                return Response({'error': 1})

            # Fetch the edited document from the Document Server. A failure here is the
            # most common cause of "edits disappear after refresh": the DS could not be
            # reached, or Django could not reach the DS's download URL.
            try:
                with urllib.request.urlopen(download_url, timeout=60) as resp:
                    content = resp.read()
            except Exception as exc:
                logger.error("OnlyOffice callback: failed to download edited doc %s from %s: %s",
                             pk, download_url, exc)
                return Response({'error': 1})

            # Never overwrite a good file with an empty/short payload.
            if not content or len(content) < 64:
                logger.error("OnlyOffice callback: refusing to save doc %s — content was "
                             "empty/too small (%d bytes)", pk, len(content or b''))
                return Response({'error': 1})

            fname = doc.file.name.rsplit('/', 1)[-1]
            try:
                from apps.core.watermark import embed_watermark
                metadata = {
                    "document_id": doc.id,
                    "title": doc.title,
                    "project_id": doc.project_id,
                    "project_number": doc.project.project_number if doc.project else None,
                    "updated_by": "onlyoffice",
                }
                content = embed_watermark(content, fname, doc.mime_type, metadata)
            except Exception as wexc:
                logger.warning("OnlyOffice callback: watermark failed for doc %s (saving "
                               "un-watermarked): %s", pk, wexc)

            try:
                doc.file.save(fname, ContentFile(content), save=False)
                doc.version += 1
                doc.save(update_fields=['file', 'version'])
                logger.info("OnlyOffice persisted doc %s -> version %s (%d bytes)",
                            pk, doc.version, len(content))
            except Exception as exc:
                logger.error("OnlyOffice callback: failed to persist doc %s: %s", pk, exc)
                return Response({'error': 1})

        return Response({'error': 0})

    @action(detail=True, methods=['post'])
    def process_ai(self, request, pk=None):
        """Queue async AI text extraction + summarisation for this document."""
        document = self.get_object()

        from apps.ai_assistant.tasks import process_document_ai
        from apps.ai_assistant.models import AITask

        # Reset flag so task re-runs even if already processed
        document.ai_processed = False
        document.save(update_fields=['ai_processed'])

        task = AITask.objects.create(
            task_type=AITask.PDF_EXTRACTION,
            requested_by=request.user,
            input_data={'document_id': document.id},
        )
        process_document_ai.delay(task.id)
        return Response({'task_id': task.id, 'detail': 'AI processing queued.'})

    @action(detail=False, methods=['post'], url_path='create-blank',
            permission_classes=[permissions.IsAuthenticated])
    def create_blank(self, request):
        """
        POST /api/documents/create-blank/
        Create a blank Word/Excel/PowerPoint document and return its doc_id for
        immediate opening in OnlyOffice.

        Body: { folder: <int>, title: <str>, doc_type: 'docx'|'xlsx'|'pptx' }
        """
        import io
        from apps.survey_projects.models import ProjectLayerFolder, SurveyProject

        folder_id = request.data.get('folder')
        title = (request.data.get('title') or 'New Document').strip()
        doc_type = request.data.get('doc_type', 'docx').lower()
        if doc_type not in ('docx', 'xlsx', 'pptx'):
            doc_type = 'docx'

        folder = None
        project = None
        if folder_id:
            try:
                folder = ProjectLayerFolder.objects.select_related('project').get(pk=folder_id)
                project = folder.project
            except ProjectLayerFolder.DoesNotExist:
                return Response({'detail': 'Folder not found.'}, status=404)

        import os
        template_name = f'new.{doc_type}'
        template_path = os.path.join(str(settings.BASE_DIR), 'apps', 'documents', 'templates', template_name)

        buf_content = None
        if os.path.exists(template_path):
            try:
                with open(template_path, 'rb') as f:
                    buf_content = f.read()
            except Exception:
                pass

        if buf_content is None:
            # Fallback to dynamic creation if template is missing
            buf = io.BytesIO()
            try:
                if doc_type == 'docx':
                    from docx import Document as DocxDoc
                    d = DocxDoc()
                    d.add_heading(title, level=0)
                    d.add_paragraph('')
                    d.save(buf)
                elif doc_type == 'xlsx':
                    import openpyxl
                    wb = openpyxl.Workbook()
                    wb.active.title = 'Sheet1'
                    wb.save(buf)
                elif doc_type == 'pptx':
                    from pptx import Presentation
                    prs = Presentation()
                    slide = prs.slides.add_slide(prs.slide_layouts[0])
                    slide.shapes.title.text = title
                    prs.save(buf)
            except ImportError as e:
                return Response({'detail': f'Required library not installed: {e}'}, status=500)
            buf.seek(0)
            buf_content = buf.read()
        safe = ''.join(c if c.isalnum() or c in '-_ ' else '_' for c in title)
        filename = f'{safe[:60]}.{doc_type}'
        file_size = len(buf_content)
        mime_map = {
            'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        }
        mime_type = mime_map.get(doc_type, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')

        doc = Document(
            project=project,
            folder=folder,
            title=title,
            category=Document.OTHER,
            file_size=file_size,
            mime_type=mime_type,
            uploaded_by=request.user,
        )
        
        # Embed watermark
        from apps.core.watermark import embed_watermark
        metadata = {
            "title": title,
            "project_id": project.id if project else None,
            "project_number": project.project_number if project else None,
            "uploaded_by": request.user.username,
        }
        try:
            buf_content = embed_watermark(buf_content, filename, mime_type, metadata)
            doc.file_size = len(buf_content)
        except Exception:
            pass
            
        doc.file.save(filename, ContentFile(buf_content), save=True)

        return Response({'doc_id': doc.id, 'title': doc.title}, status=201)

    @action(detail=False, methods=['post'], url_path='verify-watermark',
            permission_classes=[permissions.IsAuthenticated])
    def verify_watermark(self, request):
        """
        POST /api/documents/verify-watermark/
        Upload any file to verify if it was generated by RakshaGIS/DEMAP.
        """
        uploaded_file = request.FILES.get('file')
        if not uploaded_file:
            return Response({'detail': 'No file uploaded.'}, status=status.HTTP_400_BAD_REQUEST)
            
        # Secure boundary: enforce 100MB file limit
        if uploaded_file.size > 100 * 1024 * 1024:
            return Response({'detail': 'File size exceeds the 100MB limit.'}, status=status.HTTP_400_BAD_REQUEST)
            
        from apps.core.watermark import detect_watermark
        try:
            content = uploaded_file.read()
            filename = uploaded_file.name
            mime_type = getattr(uploaded_file, 'content_type', None)
            
            result = detect_watermark(content, filename, mime_type)
            return Response(result, status=status.HTTP_200_OK)
        except Exception as e:
            return Response({'detail': f'Error verifying file watermark: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
